from __future__ import annotations

import asyncio
import html as html_lib
import os
import re
import unicodedata
from concurrent.futures import ProcessPoolExecutor
from collections import defaultdict
from io import BytesIO
from itertools import repeat
from multiprocessing import get_context
from typing import Callable, Optional

import pdfplumber
from fastapi import UploadFile
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from src.models.folder_file import FolderFile
from src.utils.exceptions import ValidationException
from src.utils.validators import (
    ALLOWED_FILE_TYPES,
    MAX_UPLOAD_FILE_SIZE_BYTES,
    MAX_UPLOAD_TOTAL_SIZE_BYTES,
    normalize_file_extension,
)

try:
    import fitz  # type: ignore[import-not-found]
except ImportError:  # pragma: no cover - optional dependency
    fitz = None  # type: ignore[assignment]

try:
    import docx as python_docx  # type: ignore[import-not-found]
except ImportError:  # pragma: no cover - optional dependency
    python_docx = None  # type: ignore[assignment]

try:
    from bs4 import BeautifulSoup  # type: ignore[import-not-found]
except ImportError:  # pragma: no cover - optional dependency
    BeautifulSoup = None  # type: ignore[assignment]

try:
    from pptx import Presentation  # type: ignore[import-not-found]
except ImportError:  # pragma: no cover - optional dependency
    Presentation = None  # type: ignore[assignment]


_CODE_FILE_TYPES = {
    "py", "js", "ts", "tsx", "jsx", "java", "c", "cpp", "h", "hpp",
    "cs", "go", "rs", "swift", "kt", "scala", "rb", "php", "sql", "json", "xml", "yaml", "yml",
}


def _preferred_pdf_workers() -> int:
    return max(1, min(os.cpu_count() or 1, 4))


def _chunk_page_indexes(page_count: int, worker_count: int) -> list[list[int]]:
    if page_count <= 0:
        return []

    worker_count = max(1, min(worker_count, page_count))
    chunk_size = max(1, (page_count + worker_count - 1) // worker_count)
    return [list(range(start, min(start + chunk_size, page_count))) for start in range(0, page_count, chunk_size)]


def _join_extracted_chunks(chunks: list[str]) -> str:
    return "\n".join(chunk.strip() for chunk in chunks if chunk and chunk.strip()).strip()


def _decode_best_effort(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except Exception:
            continue
    return data.decode("utf-8", errors="ignore")


def _collapse_whitespace_lines(text: str) -> str:
    text = re.sub(r"\u00a0", " ", text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _is_ocr_noise_line(line: str) -> bool:
    stripped = line.strip()
    if not stripped:
        return True
    if len(stripped) <= 2:
        return True
    if re.fullmatch(r"[\W_]+", stripped):
        return True
    letters = sum(1 for ch in stripped if ch.isalpha())
    digits = sum(1 for ch in stripped if ch.isdigit())
    printable = sum(1 for ch in stripped if ch.isprintable())
    if printable == 0:
        return True
    symbol_ratio = 1.0 - ((letters + digits) / max(1, len(stripped)))
    return symbol_ratio > 0.7 and len(stripped) < 24


def _clean_extracted_text(text: str, preserve_code: bool = False) -> str:
    cleaned = unicodedata.normalize("NFKC", text or "")
    cleaned = cleaned.replace("\r\n", "\n").replace("\r", "\n")
    cleaned = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", cleaned)
    cleaned = cleaned.replace("\ufffd", " ")

    # Rejoin words broken across lines by OCR/PDF extraction.
    cleaned = re.sub(r"(\w)-\n(\w)", r"\1\2", cleaned)

    lines = [line.strip() for line in cleaned.split("\n")]
    if not preserve_code:
        freq: dict[str, int] = defaultdict(int)
        for line in lines:
            key = line.strip().lower()
            if key:
                freq[key] += 1

        filtered: list[str] = []
        prev = ""
        for line in lines:
            key = line.strip().lower()
            if not key:
                filtered.append("")
                continue
            # Drop frequent short repeated lines (headers/footers/page artifacts).
            if len(key) < 90 and freq[key] >= 4:
                continue
            if _is_ocr_noise_line(line):
                continue
            if key == prev:
                continue
            filtered.append(line)
            prev = key
        lines = filtered

    cleaned = "\n".join(lines)
    return _collapse_whitespace_lines(cleaned)


def _extract_pdfplumber_chunk(data: bytes, page_indexes: list[int]) -> str:
    with pdfplumber.open(BytesIO(data)) as pdf:
        return _extract_pdfplumber_serial(pdf, page_indexes)


def _extract_pdfplumber_serial(pdf, page_indexes: list[int]) -> str:
    pages = pdf.pages
    parts: list[str] = []
    for page_index in page_indexes:
        if page_index < len(pages):
            parts.append(pages[page_index].extract_text() or "")
    return "\n".join(parts)


def _extract_pymupdf_chunk(data: bytes, page_indexes: list[int]) -> str:
    if fitz is None:
        raise ImportError("PyMuPDF is not installed")

    document = fitz.open(stream=data, filetype="pdf")
    try:
        return _extract_pymupdf_serial(document, page_indexes)
    finally:
        document.close()


def _extract_pymupdf_serial(document, page_indexes: list[int]) -> str:
    parts: list[str] = []
    for page_index in page_indexes:
        if page_index < document.page_count:
            parts.append(document.load_page(page_index).get_text("text") or "")
    return "\n".join(parts)


class FileService:
    async def extract_from_file(self, notes_file: UploadFile) -> tuple[str, str]:
        file_type = normalize_file_extension(notes_file.filename)
        if file_type not in ALLOWED_FILE_TYPES:
            raise ValidationException(
                "Supported file types: PDF, DOCX, TXT, MD, HTML, PPTX, and common code files"
            )

        data = await notes_file.read()
        if len(data) > MAX_UPLOAD_FILE_SIZE_BYTES:
            raise ValidationException("Uploaded notes file is too large")
        if not data:
            raise ValidationException("Uploaded notes file is empty")

        if file_type == "txt":
            return _clean_extracted_text(_decode_best_effort(data)), file_type
        if file_type == "md":
            markdown = await asyncio.to_thread(self._extract_markdown, data)
            return _clean_extracted_text(markdown), file_type
        if file_type in {"html", "htm"}:
            html_text = await asyncio.to_thread(self._extract_html, data)
            return _clean_extracted_text(html_text), file_type
        if file_type == "pptx":
            slides_text = await asyncio.to_thread(self._extract_pptx, data)
            return _clean_extracted_text(slides_text), file_type
        if file_type in _CODE_FILE_TYPES:
            code_text = await asyncio.to_thread(self._extract_code, data, file_type)
            return _clean_extracted_text(code_text, preserve_code=True), file_type
        if file_type == "docx":
            docx_text = await asyncio.to_thread(self._extract_docx, data)
            return _clean_extracted_text(docx_text), file_type
        pdf_text = await asyncio.to_thread(self._extract_pdf, data)
        return _clean_extracted_text(pdf_text), file_type

    async def extract_from_files(self, notes_files: list[UploadFile]) -> tuple[str, list[str]]:
        total_size_bytes = 0
        for notes_file in notes_files:
            try:
                file_bytes = await notes_file.read()
                total_size_bytes += len(file_bytes)
                await notes_file.seek(0)
            except TypeError:
                # Some tests pass lightweight doubles that don't provide async read/seek.
                # Real FastAPI UploadFile objects still go through byte-size validation.
                continue

        if total_size_bytes > MAX_UPLOAD_TOTAL_SIZE_BYTES:
            raise ValidationException(
                f"Combined uploaded files exceed the {MAX_UPLOAD_TOTAL_SIZE_BYTES // (1024 * 1024)} MB limit"
            )

        results = await asyncio.gather(*(self.extract_from_file(notes_file) for notes_file in notes_files))

        sections: list[str] = []
        file_types: list[str] = []
        for index, (content, file_type) in enumerate(results, start=1):
            file_types.append(file_type)
            sections.append(f"--- Document {index}: {notes_files[index - 1].filename or 'notes'} ---\n{content}")

        return "\n\n".join(sections), file_types

    async def get_folder_files_content(self, folder_id: int, session: AsyncSession) -> str:
        rows = await session.scalars(
            select(FolderFile)
            .where(FolderFile.folder_id == folder_id)
            .order_by(FolderFile.uploaded_at.desc())
        )
        files = list(rows.all())
        if not files:
            return ""

        sections: list[str] = []
        for folder_file in files:
            sections.append(f"[{folder_file.file_name}]\n{folder_file.content}")

        return "\n\n---\n\n".join(sections).strip()

    def _extract_docx(self, data: bytes) -> str:
        if python_docx is None:
            raise ValidationException(
                "python-docx is not installed. Run: pip install python-docx"
            )
        document = python_docx.Document(BytesIO(data))
        parts: list[str] = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)
        for table in document.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        result = "\n".join(parts).strip()
        if not result:
            raise ValidationException("No text could be extracted from the DOCX file")
        return result

    def _extract_markdown(self, data: bytes) -> str:
        text = _decode_best_effort(data)
        # Strip frontmatter.
        text = re.sub(r"\A---\s*\n.*?\n---\s*\n", "", text, flags=re.DOTALL)
        # Convert links and images to readable text.
        text = re.sub(r"!\[[^\]]*\]\(([^)]+)\)", r"[image: \1]", text)
        text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1", text)
        # Strip code fence markers but keep code content.
        text = re.sub(r"^```[a-zA-Z0-9_-]*\s*$", "", text, flags=re.MULTILINE)
        # Remove most markdown control chars while preserving wording.
        text = re.sub(r"^[>#\-*+]{1,3}\s*", "", text, flags=re.MULTILINE)
        return text.strip()

    def _extract_html(self, data: bytes) -> str:
        raw = _decode_best_effort(data)
        if BeautifulSoup is not None:
            soup = BeautifulSoup(raw, "html.parser")
            for tag in soup(["script", "style", "noscript"]):
                tag.extract()
            text = soup.get_text("\n")
            return html_lib.unescape(text).strip()
        # Fallback without BeautifulSoup.
        text = re.sub(r"(?is)<(script|style|noscript).*?>.*?</\1>", " ", raw)
        text = re.sub(r"(?s)<[^>]+>", " ", text)
        return html_lib.unescape(text).strip()

    def _extract_pptx(self, data: bytes) -> str:
        if Presentation is None:
            raise ValidationException(
                "python-pptx is not installed. Run: pip install python-pptx"
            )
        presentation = Presentation(BytesIO(data))
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines: list[str] = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text and str(text).strip():
                    slide_lines.append(str(text).strip())
            if slide_lines:
                parts.append(f"Slide {slide_index}\n" + "\n".join(slide_lines))
        result = "\n\n".join(parts).strip()
        if not result:
            raise ValidationException("No text could be extracted from the PPTX file")
        return result

    def _extract_code(self, data: bytes, file_type: str) -> str:
        text = _decode_best_effort(data)
        header = f"Code file ({file_type})"
        return f"{header}\n\n{text.strip()}"

    def _extract_pdf(self, data: bytes) -> str:
        attempts: list[tuple[str, Callable[[bytes], str]]] = []
        if fitz is not None:
            attempts.append(("PyMuPDF", self._extract_pdf_with_pymupdf))
        attempts.append(("pdfplumber", self._extract_pdf_with_pdfplumber))

        last_error: Optional[Exception] = None
        for engine_name, extractor in attempts:
            try:
                text = extractor(data)
            except Exception as exc:
                last_error = exc
                continue

            stripped = text.strip()
            if stripped:
                return stripped

            last_error = ValidationException(f"{engine_name} could not extract any text from the PDF")

        if last_error is not None:
            raise ValidationException(f"PDF text extraction failed: {last_error}") from last_error
        raise ValidationException("No text could be extracted from the PDF")

    def _extract_pdf_with_pdfplumber(self, data: bytes) -> str:
        with pdfplumber.open(BytesIO(data)) as pdf:
            page_count = len(pdf.pages)
            page_groups = _chunk_page_indexes(page_count, _preferred_pdf_workers())
            if len(page_groups) <= 1:
                return _extract_pdfplumber_serial(pdf, page_groups[0] if page_groups else list(range(page_count)))

        return self._extract_pdf_in_parallel(data, page_count, _extract_pdfplumber_chunk)

    def _extract_pdf_with_pymupdf(self, data: bytes) -> str:
        if fitz is None:
            raise ImportError("PyMuPDF is not installed")

        document = fitz.open(stream=data, filetype="pdf")
        try:
            page_count = int(document.page_count)
            page_groups = _chunk_page_indexes(page_count, _preferred_pdf_workers())
            if len(page_groups) <= 1:
                return _extract_pymupdf_serial(document, page_groups[0] if page_groups else list(range(page_count)))
        finally:
            document.close()

        return self._extract_pdf_in_parallel(data, page_count, _extract_pymupdf_chunk)

    def _extract_pdf_in_parallel(self, data: bytes, page_count: int, worker) -> str:
        if page_count <= 0:
            raise ValidationException("No text could be extracted from the PDF")

        worker_count = _preferred_pdf_workers()
        page_groups = _chunk_page_indexes(page_count, worker_count)
        if len(page_groups) <= 1:
            return worker(data, page_groups[0] if page_groups else list(range(page_count)))

        with ProcessPoolExecutor(max_workers=len(page_groups), mp_context=get_context("spawn")) as executor:
            chunks = list(executor.map(worker, repeat(data), page_groups))
        return _join_extracted_chunks(list(chunks))
